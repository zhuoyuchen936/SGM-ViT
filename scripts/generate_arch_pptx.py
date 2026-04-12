#!/usr/bin/env python3
"""
scripts/generate_arch_pptx.py
==============================
ISSCC-style hardware architecture and dataflow PPTX diagrams.
White background, grayscale fills, black borders, sans-serif font.

Output: paper/EdgeStereoDAv2_architecture.pptx
  Slide 1: System Architecture
  Slide 2: Dataflow & Memory Hierarchy

Usage: python scripts/generate_arch_pptx.py
"""
import os

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

_PROJECT_DIR = os.path.normpath(os.path.join(os.path.dirname(__file__), ".."))

# ---------------------------------------------------------------------------
# ISSCC grayscale palette
# ---------------------------------------------------------------------------
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
GRAY10  = RGBColor(0x1A, 0x1A, 0x1A)   # near-black (title)
GRAY30  = RGBColor(0x4D, 0x4D, 0x4D)   # dark gray (DRAM)
GRAY50  = RGBColor(0x80, 0x80, 0x80)   # mid gray (L2 SRAM)
GRAY70  = RGBColor(0xB3, 0xB3, 0xB3)   # light gray (L1/secondary)
GRAY88  = RGBColor(0xE0, 0xE0, 0xE0)   # very light (background panels)
GRAY95  = RGBColor(0xF2, 0xF2, 0xF2)   # near-white (slide bg)

# Block fills (distinct but all grayscale)
F_VFE   = RGBColor(0xCC, 0xCC, 0xCC)   # light gray
F_CRU   = RGBColor(0xA0, 0xA0, 0xA0)   # medium-light
F_CSFE  = RGBColor(0xD8, 0xD8, 0xD8)   # lighter
F_ADCU  = RGBColor(0x8C, 0x8C, 0x8C)   # medium
F_CTRL  = RGBColor(0xBB, 0xBB, 0xBB)
F_L2    = RGBColor(0x70, 0x70, 0x70)
F_DRAM  = RGBColor(0x44, 0x44, 0x44)


def box(slide, x, y, w, h, title, fill, title_size=9, body=None, body_size=7,
        title_bold=True, title_color=BLACK, body_color=BLACK, border=BLACK):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.0)

    tf = shape.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(title_size)
    r.font.bold = title_bold
    r.font.color.rgb = title_color
    r.font.name = "Arial"

    if body:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(body_size)
        r2.font.bold = False
        r2.font.color.rgb = body_color
        r2.font.name = "Arial"
    return shape


def label(slide, x, y, w, h, text, size=8, color=BLACK, bold=False,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Arial"
    return tb


def arrow(slide, x1, y1, x2, y2, lbl=None, lbl_size=7):
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = BLACK
    conn.line.width = Pt(1.2)
    ln = conn.line._ln
    etree.SubElement(ln, qn('a:tailEnd')).set('type', 'none')
    he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'arrow'); he.set('w', 'med'); he.set('len', 'med')
    if lbl:
        mx, my = (x1 + x2) / 2, (y1 + y2) / 2
        label(slide, mx - 0.45, my - 0.15, 0.9, 0.25, lbl,
              size=lbl_size, align=PP_ALIGN.CENTER, italic=True)
    return conn


def hline(slide, x1, y, x2):
    """Thin horizontal separator line."""
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    conn.line.color.rgb = GRAY50
    conn.line.width = Pt(0.5)


# ---------------------------------------------------------------------------
# Slide 1: System Architecture
# ---------------------------------------------------------------------------
def slide1_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Title
    label(slide, 0.25, 0.05, 12.8, 0.4,
          "EdgeStereoDAv2: System Architecture",
          size=18, bold=True, color=GRAY10, align=PP_ALIGN.CENTER)
    hline(slide, 0.25, 0.5, 13.08)

    # ---- DRAM bar (bottom) ----
    box(slide, 0.25, 6.55, 12.83, 0.5,
        "L3 DRAM  (LPDDR4x · 2 GB · 16 GB/s)",
        F_DRAM, title_size=9, title_bold=True, title_color=WHITE)

    # ---- L2 SRAM bar ----
    box(slide, 0.25, 5.75, 12.83, 0.6,
        "L2 Global Shared SRAM",
        F_L2, title_size=9, title_bold=True, title_color=WHITE,
        body="512 KB · 32 banks · 64 B/cycle · 4-ch DMA · Priority arbiter: VFE > CSFE > ADCU > DMA",
        body_size=7.5, body_color=WHITE)

    # ---- Compute blocks ----
    # VFE
    box(slide, 0.25, 1.85, 3.6, 3.7,
        "VFE",
        F_VFE, title_size=14, title_bold=True,
        body="32×32 INT8 Systolic Array\n1024 MACs · WS Dataflow\nFlash Attn (Br=64, Bc=128)\nSoftmax: 16-seg PWL\nGELU: 32-seg PWL\nLayerNorm: 3-stage pipeline\nL1 SRAM: 64 KB",
        body_size=8)

    # CRU
    box(slide, 4.1, 1.85, 2.0, 3.7,
        "CRU",
        F_CRU, title_size=14, title_bold=True, title_color=WHITE,
        body="1369 Comparators\nPrefix-Sum Index Gen\nIndex SRAM: 2.7 KB\nLatency: 1 cycle\nArea: <0.3%",
        body_size=8, body_color=WHITE)

    # CSFE
    box(slide, 6.35, 1.85, 3.2, 3.7,
        "CSFE",
        F_CSFE, title_size=14, title_bold=True,
        body="16×16 MACs (256) · OS Mode\n1×1 + 3×3 Conv\nBilinear Upsample: 32 px/cyc\nStages: 37→74→148→296\nL1 SRAM: 64 KB",
        body_size=8)

    # ADCU
    box(slide, 9.8, 1.85, 2.9, 3.7,
        "ADCU",
        F_ADCU, title_size=14, title_bold=True, title_color=WHITE,
        body="Sparse NCC Matcher\n32 kpts · 11×11 patch\n2×2 LS Solver (106 cyc)\nReciprocal LUT: 4096 entries\nZ=αd+β → 1/Z → d=Bf/Z\n32 px/cycle · L1: 8 KB",
        body_size=8, body_color=WHITE)

    # CTRL (narrow right)
    box(slide, 12.9, 1.85, 0.93, 3.7,
        "CTRL",
        F_CTRL, title_size=9, title_bold=True,
        body="RISC-V\nRV32I\n16 KB\nIROM\n4 KB\nDRAM",
        body_size=7)

    # ---- I/O labels ----
    label(slide, 0.25, 1.55, 3.6, 0.28,
          "Input: Left image (from DRAM)", size=7.5, italic=True, color=GRAY30)
    label(slide, 9.8, 1.55, 2.9, 0.28,
          "Input: Stereo L+R (from DRAM)", size=7.5, italic=True, color=GRAY30)
    label(slide, 9.8, 5.55, 2.9, 0.2,
          "Output: Disparity map → DRAM", size=7.5, italic=True, color=GRAY30)

    # ---- Data flow arrows ----
    # ADCU → CRU: confidence map
    arrow(slide, 9.8, 3.75, 6.1, 3.75, lbl="conf map + α,β")
    # CRU → VFE: keep indices
    arrow(slide, 4.1, 3.75, 3.85, 3.75, lbl="keep_idx")
    # VFE → CSFE: DPT features
    arrow(slide, 3.85, 3.75, 6.35, 3.75, lbl="feat [2,5,8,11]")
    # CSFE → ADCU: relative depth
    arrow(slide, 9.55, 3.75, 9.8, 3.75, lbl="d_rel")

    # L2 ↔ blocks (vertical)
    arrow(slide, 2.05, 5.55, 2.05, 5.75)
    arrow(slide, 7.95, 5.55, 7.95, 5.75)
    arrow(slide, 11.25, 5.55, 11.25, 5.75)
    # L2 ↔ DRAM
    arrow(slide, 6.67, 6.35, 6.67, 6.55)

    # ---- Abbreviation footnote ----
    hline(slide, 0.25, 7.1, 13.08)
    abbrevs = ("VFE: ViT Feature Engine  |  CRU: Confidence Router Unit  |  "
               "CSFE: Cross-Scale Fusion Engine  |  ADCU: Absolute Disparity Calibration Unit  |  "
               "CTRL: Control Processor  |  WS: Weight-Stationary  |  OS: Output-Stationary  |  "
               "GAS: Gather-Attend-Scatter  |  PWL: Piecewise-Linear  |  NCC: Normalized Cross-Correlation  |  "
               "LS: Least-Squares  |  LUT: Look-Up Table  |  DMA: Direct Memory Access")
    label(slide, 0.25, 7.12, 12.83, 0.35, abbrevs,
          size=6.5, color=GRAY30, italic=True)


# ---------------------------------------------------------------------------
# Slide 2: Dataflow & Memory Hierarchy
# ---------------------------------------------------------------------------
def slide2_dataflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    label(slide, 0.25, 0.05, 12.8, 0.4,
          "EdgeStereoDAv2: Dataflow & Memory Hierarchy",
          size=18, bold=True, color=GRAY10, align=PP_ALIGN.CENTER)
    hline(slide, 0.25, 0.5, 13.08)

    # ---- Pipeline stages ----
    stages = [
        ("ADCU\n(Sparse Match)",
         "32 kpts · 11×11 NCC\n512 cyc matching\n+106 cyc LS solve",
         F_ADCU, WHITE, 0.25),
        ("CRU\n(Token Route)",
         "1369 comparators\nθ threshold\nprefix-sum compaction\n1 cycle total",
         F_CRU, WHITE, 2.85),
        ("VFE\n(ViT Encoder)",
         "12 transformer blocks\nWS 32×32 systolic\nGAS sparse attention\nFlash Attn: Br=64, Bc=128",
         F_VFE, BLACK, 5.45),
        ("CSFE\n(DPT Decoder)",
         "4 fusion stages\nOS 16×16 conv array\n37→74→148→296 px\n3×3 conv + 2× upsample",
         F_CSFE, BLACK, 8.35),
        ("ADCU\n(Depth→Disp)",
         "Z = αd_rel + β\n1/Z via 4096-entry LUT\nd = Bf/Z\n32 px/cycle",
         F_ADCU, WHITE, 11.05),
    ]

    for title, body, fill, tc, x in stages:
        box(slide, x, 0.6, 2.35, 2.3, title, fill,
            title_size=9, title_bold=True, title_color=tc,
            body=body, body_size=7.5, body_color=tc)

    # Arrows + labels between stages
    flow_labels = ["conf map\n+ α, β", "keep_mask\nkeep_idx",
                   "feat\n[2,5,8,11]", "d_rel"]
    for i, lbl in enumerate(flow_labels):
        x1 = stages[i][4] + 2.35
        x2 = stages[i+1][4]
        arrow(slide, x1, 1.75, x2, 1.75, lbl=lbl)

    # Input/output
    label(slide, 0.25, 0.35, 2.35, 0.25,
          "← Stereo L+R (DRAM)", size=7, italic=True, color=GRAY30)
    label(slide, 11.05, 0.35, 2.35, 0.25,
          "→ Disparity (DRAM)", size=7, italic=True, color=GRAY30, align=PP_ALIGN.RIGHT)

    # ---- Memory hierarchy (left column) ----
    label(slide, 0.25, 3.1, 5.5, 0.3,
          "Memory Hierarchy", size=10, bold=True, color=GRAY10)

    mem = [
        ("L0  Register File",
         "512 B/PE · 64 B/cyc · 0 lat · 0.05 pJ/acc  —  MAC operands, partial sums, flash-attn running stats",
         GRAY88, BLACK),
        ("L1  Local SRAM  (per engine)",
         "64 KB · 32 B/cyc · 1 cyc lat · 2.0 pJ/acc · 16 banks  —  Weight tile, activation tile, attn score tile",
         GRAY70, BLACK),
        ("L2  Global Shared SRAM",
         "512 KB · 64 B/cyc · 3 cyc lat · 10.0 pJ/acc · 32 banks  —  Encoder features, weight prefetch (128 KB), feature buf (256 KB)",
         GRAY50, WHITE),
        ("L3  DRAM  (LPDDR4x)",
         "2 GB · 4 B/cyc · 50 cyc lat · 200 pJ/acc  —  Full model weights (~25 MB INT8), frame buffers",
         GRAY30, WHITE),
    ]
    for i, (name, desc, fill, tc) in enumerate(mem):
        box(slide, 0.25, 3.45 + i * 0.62, 6.5, 0.55,
            name, fill, title_size=8.5, title_bold=True, title_color=tc,
            body=desc, body_size=7, body_color=tc)

    # ---- Dataflow detail (right column) ----
    label(slide, 7.1, 3.1, 6.0, 0.3,
          "Dataflow Modes", size=10, bold=True, color=GRAY10)

    ws_body = ("Weights preloaded into PE registers (stationary)\n"
               "Activations flow left→right · Partial sums top→bottom\n"
               "Tile: M=64, N=128, K=64\n"
               "Ops: QKV proj · Q·Kᵀ · Attn·V · Out proj · MLP W1/W2\n"
               "Arith. intensity: 128–502 OPS/Byte  →  compute-bound")
    box(slide, 7.1, 3.45, 6.0, 1.5,
        "Weight-Stationary (WS)  —  VFE Encoder",
        F_VFE, title_size=9, title_bold=True,
        body=ws_body, body_size=7.5)

    os_body = ("Output pixels remain in PE accumulators\n"
               "Weights and inputs stream through\n"
               "Tile: H=8, W=8, Co=16, Ci=32\n"
               "Ops: 1×1 proj · 3×3 RCU conv · bilinear 2× upsample\n"
               "256 MACs (16×16 subset) · ~80% utilization\n"
               "Reconfig: 1 cycle between 1×1 and 3×3 modes")
    box(slide, 7.1, 5.1, 6.0, 1.6,
        "Output-Stationary (OS)  —  CSFE Decoder",
        F_CSFE, title_size=9, title_bold=True,
        body=os_body, body_size=7.5)

    # ---- Abbreviation footnote ----
    hline(slide, 0.25, 7.1, 13.08)
    abbrevs = ("VFE: ViT Feature Engine  |  CRU: Confidence Router Unit  |  "
               "CSFE: Cross-Scale Fusion Engine  |  ADCU: Absolute Disparity Calibration Unit  |  "
               "WS: Weight-Stationary  |  OS: Output-Stationary  |  GAS: Gather-Attend-Scatter  |  "
               "DPT: Dense Prediction Transformer  |  NCC: Normalized Cross-Correlation  |  "
               "LS: Least-Squares  |  LUT: Look-Up Table  |  PWL: Piecewise-Linear  |  "
               "kpts: keypoints  |  cyc: cycles  |  lat: latency  |  acc: access  |  buf: buffer")
    label(slide, 0.25, 7.12, 12.83, 0.35, abbrevs,
          size=6.5, color=GRAY30, italic=True)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide1_architecture(prs)
    slide2_dataflow(prs)
    out = os.path.join(_PROJECT_DIR, 'paper', 'EdgeStereoDAv2_architecture.pptx')
    prs.save(out)
    print(f'Saved: {out}')


if __name__ == '__main__':
    main()
